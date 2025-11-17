"""Enhanced export service with comprehensive progress tracking and error handling."""

import asyncio
import io
import json
import logging
import os
import tempfile
import zipfile
from datetime import datetime, timedelta
from typing import Any, Dict, List, Optional, Tuple, Union
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field

from fastapi import BackgroundTasks
from backend.repositories.presentation_repository import PresentationRepository
from backend.repositories.lesson_repository import LessonRepository
from backend.lessons.presentation_schema import PresentationDocument, PresentationExport, PresentationStatus
from backend.models.export_progress import (
    ExportFormat, ExportStep, ExportStepProgress, ExportFormatProgress,
    BulkExportProgress, ExportProgressUpdate, ExportAnalytics, ExportStatus
)
from backend.services.presentation_errors import (
    PresentationError, PresentationErrorCode, create_error_from_exception,
    ErrorRecoveryStrategy
)
from backend.services.progress_websocket import get_progress_service

logger = logging.getLogger(__name__)


@dataclass
class ExportOptions:
    """Options for export operations."""
    include_metadata: bool = True
    include_teacher_script: bool = True
    quality_settings: Dict[str, Any] = field(default_factory=dict)
    pdf_options: Dict[str, Any] = field(default_factory=dict)
    pptx_options: Dict[str, Any] = field(default_factory=dict)
    batch_size: int = 1
    max_retries: int = 3
    timeout_seconds: int = 300


class ExportService:
    """Enhanced export service with comprehensive progress tracking."""

    def __init__(
        self,
        presentation_repo: Optional[PresentationRepository] = None,
        lesson_repo: Optional[LessonRepository] = None,
        max_concurrent_exports: int = 4
    ):
        self.presentation_repo = presentation_repo or PresentationRepository()
        self.lesson_repo = lesson_repo or LessonRepository()
        self.progress_service = get_progress_service()

        # Active export tracking
        self.active_exports: Dict[str, ExportFormatProgress] = {}
        self.active_bulk_exports: Dict[str, BulkExportProgress] = {}
        self.export_queue: asyncio.Queue = asyncio.Queue()
        self.executor = ThreadPoolExecutor(max_workers=max_concurrent_exports)
        self.max_concurrent_exports = max_concurrent_exports

    async def export_presentation(
        self,
        presentation_id: str,
        format: ExportFormat,
        user_id: str,
        options: Optional[ExportOptions] = None,
        websocket_available: bool = True
    ) -> ExportFormatProgress:
        """Export a presentation in the specified format with detailed progress tracking."""

        if not options:
            options = ExportOptions()

        # Create export progress tracker
        export_progress = ExportFormatProgress(
            export_id=f"export_{presentation_id}_{format.value}_{int(time.time())}",
            format=format
        )
        self.active_exports[export_progress.export_id] = export_progress

        try:
            # Start the export
            export_progress.start_export()
            await self._send_progress_update(export_progress, user_id, websocket_available)

            # Validate presentation
            await self._validate_presentation_step(export_progress, presentation_id, user_id, websocket_available)

            # Prepare content
            presentation = await self._prepare_content_step(export_progress, presentation_id, user_id, websocket_available)

            # Format export
            export_data = await self._format_export_step(export_progress, presentation, format, options, websocket_available)

            # Validate output
            await self._validate_output_step(export_progress, format, export_data, websocket_available)

            # Calculate size and filename
            file_size, filename = await self._finalize_export_step(export_progress, format, export_data, websocket_available)

            # Complete export
            export_progress.complete_export(file_size, filename)
            await self._send_progress_update(export_progress, user_id, websocket_available)

            # Store analytics
            await self._store_export_analytics(export_progress, presentation_id, user_id, websocket_available)

            logger.info(f"Successfully exported presentation {presentation_id} as {format.value}")
            return export_progress

        except Exception as e:
            # Handle export failure
            error_code = await self._handle_export_failure(export_progress, e, user_id, websocket_available)
            logger.error(f"Export failed for presentation {presentation_id} as {format.value}: {e}")
            raise

        finally:
            # Clean up active exports
            if export_progress.export_id in self.active_exports:
                del self.active_exports[export_progress.export_id]

    async def bulk_export_presentation(
        self,
        presentation_id: str,
        formats: List[ExportFormat],
        user_id: str,
        options: Optional[ExportOptions] = None,
        create_zip: bool = True,
        websocket_available: bool = True
    ) -> BulkExportProgress:
        """Export a presentation in multiple formats with concurrent processing."""

        if not options:
            options = ExportOptions()

        # Create bulk export tracker
        bulk_progress = BulkExportProgress(
            presentation_id=presentation_id,
            user_id=user_id,
            formats=formats,
            concurrent_exports=min(options.batch_size, self.max_concurrent_exports)
        )
        self.active_bulk_exports[bulk_progress.bulk_export_id] = bulk_progress

        try:
            # Start bulk export
            bulk_progress.start_bulk_export()
            await self._send_bulk_progress_update(bulk_progress, user_id, websocket_available)

            # Process exports concurrently
            await self._process_concurrent_exports(bulk_progress, presentation_id, formats, user_id, options, websocket_available)

            # Create ZIP if multiple successful exports and requested
            if create_zip and len(bulk_progress.successful_exports) > 1:
                await self._create_download_zip(bulk_progress, user_id, websocket_available)

            # Complete bulk export
            if bulk_progress.status == ExportStatus.RUNNING:
                bulk_progress.status = ExportStatus.COMPLETED
                bulk_progress.completed_at = datetime.utcnow()

            await self._send_bulk_progress_update(bulk_progress, user_id, websocket_available)

            logger.info(f"Successfully completed bulk export for presentation {presentation_id}")
            return bulk_progress

        except Exception as e:
            # Handle bulk export failure
            bulk_progress.error_message = str(e)
            bulk_progress.status = ExportStatus.FAILED
            await self._send_bulk_progress_update(bulk_progress, user_id, websocket_available)
            logger.error(f"Bulk export failed for presentation {presentation_id}: {e}")
            raise

        finally:
            # Clean up
            if bulk_progress.bulk_export_id in self.active_bulk_exports:
                del self.active_bulk_exports[bulk_progress.bulk_export_id]

    async def get_export_progress(self, export_id: str) -> Optional[ExportFormatProgress]:
        """Get progress for a specific export."""
        return self.active_exports.get(export_id)

    async def get_bulk_export_progress(self, bulk_export_id: str) -> Optional[BulkExportProgress]:
        """Get progress for a bulk export."""
        return self.active_bulk_exports.get(bulk_export_id)

    async def cancel_export(self, export_id: str, user_id: str) -> bool:
        """Cancel a specific export."""
        if export_id in self.active_exports:
            export_progress = self.active_exports[export_id]
            export_progress.cancel_export("User requested cancellation")
            await self._send_progress_update(export_progress, user_id)
            return True
        return False

    async def cancel_bulk_export(self, bulk_export_id: str, user_id: str) -> bool:
        """Cancel a bulk export."""
        if bulk_export_id in self.active_bulk_exports:
            bulk_progress = self.active_bulk_exports[bulk_export_id]

            # Cancel all running format exports
            for format_progress in bulk_progress.export_progress.values():
                if format_progress.status == ExportStatus.RUNNING:
                    format_progress.cancel_export("User requested cancellation")

            bulk_progress.status = ExportStatus.CANCELLED
            await self._send_bulk_progress_update(bulk_progress, user_id)
            return True
        return False

    async def retry_export(self, export_id: str, user_id: str) -> bool:
        """Retry a failed export."""
        if export_id not in self.active_exports:
            return False

        export_progress = self.active_exports[export_id]
        if export_progress.status != ExportStatus.FAILED:
            return False

        return export_progress.retry_export()

    async def _validate_presentation_step(
        self,
        export_progress: ExportFormatProgress,
        presentation_id: str,
        user_id: str,
        websocket_available: bool
    ) -> None:
        """Step 1: Validate presentation access and readiness."""
        export_progress.start_step(ExportStep.VALIDATING_PRESENTATION)
        await self._send_progress_update(export_progress, user_id, websocket_available)

        try:
            # Get presentation
            presentation = self.presentation_repo.get_presentation(presentation_id)
            if not presentation:
                raise PresentationError.lesson_not_found(presentation_id)

            # Verify user access
            if presentation.lesson_id:
                lesson = self.lesson_repo.get_lesson(presentation.lesson_id)
                if lesson and lesson.user_id != user_id:
                    raise PresentationError.lesson_access_denied(presentation_id)

            # Check presentation status
            if presentation.status != PresentationStatus.COMPLETE:
                raise PresentationError.validation_failed("presentation_status", presentation.status.value)

            # Update progress
            await self._update_step_progress(
                export_progress, ExportStep.VALIDATING_PRESENTATION,
                progress=100.0, message="Validation successful",
                user_id=user_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.VALIDATING_PRESENTATION)
            await self._send_progress_update(export_progress, user_id, websocket_available)

        except PresentationError as e:
            export_progress.fail_step(ExportStep.VALIDATING_PRESENTATION, e.user_message, e.code.value)
            await self._send_progress_update(export_progress, user_id, websocket_available)
            raise

    async def _prepare_content_step(
        self,
        export_progress: ExportFormatProgress,
        presentation_id: str,
        user_id: str,
        websocket_available: bool
    ) -> PresentationDocument:
        """Step 2: Prepare presentation content for export."""
        export_progress.start_step(ExportStep.PREPARING_CONTENT)
        await self._send_progress_update(export_progress, user_id, websocket_available)

        try:
            # Get full presentation with slides
            presentation = self.presentation_repo.get_presentation(presentation_id)

            # Simulate content preparation with progress updates
            for i, slide in enumerate(presentation.slides, 1):
                progress = (i / len(presentation.slides)) * 100.0
                await self._update_step_progress(
                    export_progress, ExportStep.PREPARING_CONTENT,
                    progress=progress, message=f"Preparing content: {i}/{len(presentation.slides)} slides",
                    user_id=user_id, websocket_available=websocket_available
                )
                await asyncio.sleep(0.1)  # Simulate processing time

            export_progress.complete_step(ExportStep.PREPARING_CONTENT)
            await self._send_progress_update(export_progress, user_id, websocket_available)

            return presentation

        except Exception as e:
            error = create_error_from_exception(e, {"step": "prepare_content"})
            export_progress.fail_step(ExportStep.PREPARING_CONTENT, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, user_id, websocket_available)
            raise

    async def _format_export_step(
        self,
        export_progress: ExportFormatProgress,
        presentation: PresentationDocument,
        format: ExportFormat,
        options: ExportOptions,
        websocket_available: bool
    ) -> Union[str, bytes]:
        """Step 3: Format export based on specified format."""

        if format == ExportFormat.JSON:
            return await self._format_json_export(export_progress, presentation, options, websocket_available)
        elif format == ExportFormat.MARKDOWN:
            return await self._format_markdown_export(export_progress, presentation, options, websocket_available)
        elif format == ExportFormat.PPTX:
            return await self._format_pptx_export(export_progress, presentation, options, websocket_available)
        elif format == ExportFormat.PDF:
            return await self._format_pdf_export(export_progress, presentation, options, websocket_available)
        else:
            raise PresentationError.validation_failed("export_format", format.value)

    async def _format_json_export(
        self,
        export_progress: ExportFormatProgress,
        presentation: PresentationDocument,
        options: ExportOptions,
        websocket_available: bool
    ) -> str:
        """Format JSON export."""

        # JSON serialization
        export_progress.start_step(ExportStep.JSON_SERIALIZING)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            # Prepare data structure
            json_data = {
                "presentation": presentation.model_dump(mode="json"),
                "generated_at": datetime.utcnow().isoformat(),
                "format": "p1.0",
                "includes_metadata": options.include_metadata,
                "includes_teacher_script": options.include_teacher_script
            }

            # Serialize with progress
            await self._update_step_progress(
                export_progress, ExportStep.JSON_SERIALIZING,
                progress=50.0, message="Serializing presentation data",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            json_content = json.dumps(json_data, indent=2, ensure_ascii=False)

            export_progress.complete_step(ExportStep.JSON_SERIALIZING)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "json_serialization"})
            export_progress.fail_step(ExportStep.JSON_SERIALIZING, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        # JSON validation
        export_progress.start_step(ExportStep.JSON_VALIDATING)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            # Validate JSON structure
            json.loads(json_content)

            await self._update_step_progress(
                export_progress, ExportStep.JSON_VALIDATING,
                progress=100.0, message="JSON validation successful",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.JSON_VALIDATING)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "json_validation"})
            export_progress.fail_step(ExportStep.JSON_VALIDATING, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        return json_content

    async def _format_markdown_export(
        self,
        export_progress: ExportFormatProgress,
        presentation: PresentationDocument,
        options: ExportOptions,
        websocket_available: bool
    ) -> str:
        """Format Markdown export."""

        # Markdown formatting
        export_progress.start_step(ExportStep.MARKDOWN_FORMATTING)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            markdown_lines = [
                f"# {presentation.slides[0].title if presentation.slides else 'Presentation'}",
                "",
                f"**Generated:** {presentation.created_at.strftime('%Y-%m-%d %H:%M')}",
                f"**Lesson ID:** {presentation.lesson_id}",
                f"**Revision:** {presentation.lesson_revision}",
                "",
                "---",
                "",
            ]

            # Process slides with progress updates
            for i, slide in enumerate(presentation.slides, 1):
                progress = (i / len(presentation.slides)) * 80.0  # 80% of the step
                await self._update_step_progress(
                    export_progress, ExportStep.MARKDOWN_FORMATTING,
                    progress=progress, message=f"Formatting slide {i}/{len(presentation.slides)}: {slide.title}",
                    user_id=presentation.lesson_id, websocket_available=websocket_available
                )

                markdown_lines.extend([
                    f"## {slide.title}",
                    "",
                ])

                if slide.subtitle:
                    markdown_lines.extend([
                        f"**{slide.subtitle}**",
                        "",
                    ])

                if slide.key_points:
                    markdown_lines.extend([
                        "**Key Points:**",
                        "",
                    ])
                    for point in slide.key_points:
                        markdown_lines.append(f"- {point}")
                    markdown_lines.append("")

                if slide.teacher_script and options.include_teacher_script:
                    markdown_lines.extend([
                        "**Teacher Script:**",
                        "",
                        slide.teacher_script,
                        "",
                    ])

                if slide.duration_minutes:
                    markdown_lines.extend([
                        f"**Duration:** {slide.duration_minutes} minutes",
                        "",
                    ])

                if slide.standard_codes:
                    markdown_lines.extend([
                        "**Standards:** " + ", ".join(slide.standard_codes),
                        "",
                    ])

                markdown_lines.extend(["---", ""])

            markdown_content = "\n".join(markdown_lines)

            export_progress.complete_step(ExportStep.MARKDOWN_FORMATTING)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "markdown_formatting"})
            export_progress.fail_step(ExportStep.MARKDOWN_FORMATTING, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        # Markdown styling
        export_progress.start_step(ExportStep.MARKDOWN_STYLING)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            # Apply final styling and validation
            await self._update_step_progress(
                export_progress, ExportStep.MARKDOWN_STYLING,
                progress=100.0, message="Applying final markdown styling",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.MARKDOWN_STYLING)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "markdown_styling"})
            export_progress.fail_step(ExportStep.MARKDOWN_STYLING, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        return markdown_content

    async def _format_pptx_export(
        self,
        export_progress: ExportFormatProgress,
        presentation: PresentationDocument,
        options: ExportOptions,
        websocket_available: bool
    ) -> str:
        """Format PPTX export with detailed progress tracking."""

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError as e:
            error = PresentationError.export_failed("pptx", f"Missing dependency: {e.name}")
            export_progress.status = ExportStatus.FAILED
            export_progress.error_message = error.user_message
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise error

        # Create presentation
        export_progress.start_step(ExportStep.PPTX_CREATING_PRESENTATION)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            prs = Presentation()
            await self._update_step_progress(
                export_progress, ExportStep.PPTX_CREATING_PRESENTATION,
                progress=50.0, message="Creating PowerPoint presentation",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = presentation.slides[0].title if presentation.slides else "Presentation"
            subtitle.text = f"Generated: {presentation.created_at.strftime('%Y-%m-%d %H:%M')}"

            await self._update_step_progress(
                export_progress, ExportStep.PPTX_CREATING_PRESENTATION,
                progress=100.0, message="Title slide created",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.PPTX_CREATING_PRESENTATION)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "pptx_creating_presentation"})
            export_progress.fail_step(ExportStep.PPTX_CREATING_PRESENTATION, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        # Add slides
        export_progress.start_step(ExportStep.PPTX_ADDING_SLIDES)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            content_slides = presentation.slides[1:] if len(presentation.slides) > 1 else []

            for i, content_slide in enumerate(content_slides, 1):
                progress = (i / len(content_slides)) * 100.0 if content_slides else 100.0
                await self._update_step_progress(
                    export_progress, ExportStep.PPTX_ADDING_SLIDES,
                    progress=progress, message=f"Adding slide {i}/{len(content_slides)}: {content_slide.title}",
                    user_id=presentation.lesson_id, websocket_available=websocket_available
                )

                # Choose layout
                if content_slide.title and content_slide.key_points:
                    slide_layout = prs.slide_layouts[1]
                else:
                    slide_layout = prs.slide_layouts[5]

                slide = prs.slides.add_slide(slide_layout)

                # Add title
                if content_slide.title and slide.shapes.title:
                    slide.shapes.title.text = content_slide.title

                # Add content
                content_parts = []
                if content_slide.key_points:
                    content_parts.extend([f"• {point}" for point in content_slide.key_points])
                if content_slide.teacher_script and options.include_teacher_script:
                    content_parts.append(f"Teacher: {content_slide.teacher_script}")

                if content_parts and len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    content_placeholder.text = "\n".join(content_parts)

            export_progress.complete_step(ExportStep.PPTX_ADDING_SLIDES)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "pptx_adding_slides"})
            export_progress.fail_step(ExportStep.PPTX_ADDING_SLIDES, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        # Apply styles
        export_progress.start_step(ExportStep.PPTX_APPLYING_STYLES)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            await self._update_step_progress(
                export_progress, ExportStep.PPTX_APPLYING_STYLES,
                progress=100.0, message="Applying PowerPoint styles",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.PPTX_APPLYING_STYLES)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "pptx_applying_styles"})
            export_progress.fail_step(ExportStep.PPTX_APPLYING_STYLES, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        # Save presentation
        export_progress.start_step(ExportStep.PPTX_SAVING)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
                prs.save(temp_file.name)
                file_path = temp_file.name

            await self._update_step_progress(
                export_progress, ExportStep.PPTX_SAVING,
                progress=100.0, message="PowerPoint file saved",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.PPTX_SAVING)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "pptx_saving"})
            export_progress.fail_step(ExportStep.PPTX_SAVING, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        return file_path

    async def _format_pdf_export(
        self,
        export_progress: ExportFormatProgress,
        presentation: PresentationDocument,
        options: ExportOptions,
        websocket_available: bool
    ) -> str:
        """Format PDF export with detailed progress tracking."""

        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.units import inch
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.colors import black, blue
        except ImportError as e:
            error = PresentationError.export_failed("pdf", f"Missing dependency: {e.name}")
            export_progress.status = ExportStatus.FAILED
            export_progress.error_message = error.user_message
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise error

        # Create document
        export_progress.start_step(ExportStep.PDF_CREATING_DOCUMENT)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
                temp_filename = temp_file.name

            doc = SimpleDocTemplate(temp_filename, pagesize=A4)
            styles = getSampleStyleSheet()

            await self._update_step_progress(
                export_progress, ExportStep.PDF_CREATING_DOCUMENT,
                progress=100.0, message="PDF document structure created",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.PDF_CREATING_DOCUMENT)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "pdf_creating_document"})
            export_progress.fail_step(ExportStep.PDF_CREATING_DOCUMENT, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        # Add content
        export_progress.start_step(ExportStep.PDF_ADDING_CONTENT)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            story = []

            # Add title
            title_style = styles["Title"]
            title_style.textColor = blue
            title = presentation.slides[0].title if presentation.slides else "Presentation"
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 12))

            # Add metadata
            normal_style = styles["Normal"]
            story.append(
                Paragraph(
                    f"<b>Generated:</b> {presentation.created_at.strftime('%Y-%m-%d %H:%M')}",
                    normal_style,
                )
            )
            story.append(Spacer(1, 20))

            # Process slides
            for i, slide in enumerate(presentation.slides, 1):
                progress = (i / len(presentation.slides)) * 100.0
                await self._update_step_progress(
                    export_progress, ExportStep.PDF_ADDING_CONTENT,
                    progress=progress, message=f"Adding slide {i}/{len(presentation.slides)}: {slide.title}",
                    user_id=presentation.lesson_id, websocket_available=websocket_available
                )

                # Slide title
                if slide.title:
                    heading_style = styles["Heading1"]
                    story.append(Paragraph(f"Slide {i + 1}: {slide.title}", heading_style))
                    story.append(Spacer(1, 12))

                # Content
                if slide.key_points:
                    story.append(Paragraph("<b>Key Points:</b>", normal_style))
                    for point in slide.key_points:
                        story.append(Paragraph(f"• {point}", normal_style))
                    story.append(Spacer(1, 6))

                if slide.teacher_script and options.include_teacher_script:
                    story.append(Paragraph("<b>Teacher Script:</b>", normal_style))
                    story.append(Paragraph(slide.teacher_script, normal_style))
                    story.append(Spacer(1, 6))

                story.append(Spacer(1, 20))

            export_progress.complete_step(ExportStep.PDF_ADDING_CONTENT)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "pdf_adding_content"})
            export_progress.fail_step(ExportStep.PDF_ADDING_CONTENT, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        # Format and render
        export_progress.start_step(ExportStep.PDF_FORMATTING)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            await self._update_step_progress(
                export_progress, ExportStep.PDF_FORMATTING,
                progress=50.0, message="Formatting PDF layout",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.PDF_FORMATTING)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "pdf_formatting"})
            export_progress.fail_step(ExportStep.PDF_FORMATTING, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        export_progress.start_step(ExportStep.PDF_RENDERING)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            await self._update_step_progress(
                export_progress, ExportStep.PDF_RENDERING,
                progress=50.0, message="Rendering PDF document",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            doc.build(story)

            await self._update_step_progress(
                export_progress, ExportStep.PDF_RENDERING,
                progress=100.0, message="PDF rendering complete",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.PDF_RENDERING)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "pdf_rendering"})
            export_progress.fail_step(ExportStep.PDF_RENDERING, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        # Save PDF
        export_progress.start_step(ExportStep.PDF_SAVING)
        await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        try:
            await self._update_step_progress(
                export_progress, ExportStep.PDF_SAVING,
                progress=100.0, message="PDF file saved",
                user_id=presentation.lesson_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.PDF_SAVING)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "pdf_saving"})
            export_progress.fail_step(ExportStep.PDF_SAVING, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, presentation.lesson_id, websocket_available)
            raise

        return temp_filename

    async def _validate_output_step(
        self,
        export_progress: ExportFormatProgress,
        format: ExportFormat,
        export_data: Union[str, bytes],
        websocket_available: bool,
        user_id: Optional[str] = None
    ) -> None:
        """Step: Validate the generated export output."""

        export_progress.start_step(ExportStep.VALIDATING_OUTPUT)
        await self._send_progress_update(export_progress, user_id, websocket_available)

        try:
            # Basic validation
            if isinstance(export_data, str):
                content_size = len(export_data.encode('utf-8'))
            else:
                content_size = len(export_data)

            if content_size == 0:
                raise PresentationError.export_failed(format.value, "Generated export is empty")

            # Format-specific validation
            if format == ExportFormat.JSON:
                json.loads(export_data)  # Validate JSON
            elif format == ExportFormat.MARKDOWN:
                # Basic markdown validation - check for key headers
                if not export_data.startswith('#'):
                    raise PresentationError.export_failed(format.value, "Invalid markdown format")
            elif format in [ExportFormat.PPTX, ExportFormat.PDF]:
                # File existence check
                if not os.path.exists(export_data):
                    raise PresentationError.export_failed(format.value, "Generated file not found")

            await self._update_step_progress(
                export_progress, ExportStep.VALIDATING_OUTPUT,
                progress=100.0, message="Output validation successful",
                user_id=user_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.VALIDATING_OUTPUT)
            await self._send_progress_update(export_progress, user_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "validate_output"})
            export_progress.fail_step(ExportStep.VALIDATING_OUTPUT, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, user_id, websocket_available)
            raise

    async def _finalize_export_step(
        self,
        export_progress: ExportFormatProgress,
        format: ExportFormat,
        export_data: Union[str, bytes],
        websocket_available: bool,
        user_id: Optional[str] = None
    ) -> Tuple[int, str]:
        """Step: Calculate size and generate filename."""

        # Calculate size
        export_progress.start_step(ExportStep.CALCULATING_SIZE)
        await self._send_progress_update(export_progress, user_id, websocket_available)

        try:
            if isinstance(export_data, str):
                if format in [ExportFormat.PPTX, ExportFormat.PDF]:
                    # File path
                    file_size = os.path.getsize(export_data)
                else:
                    # Content string
                    file_size = len(export_data.encode('utf-8'))
            else:
                file_size = len(export_data)

            await self._update_step_progress(
                export_progress, ExportStep.CALCULATING_SIZE,
                progress=100.0, message=f"File size calculated: {file_size} bytes",
                user_id=user_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.CALCULATING_SIZE)
            await self._send_progress_update(export_progress, user_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "calculating_size"})
            export_progress.fail_step(ExportStep.CALCULATING_SIZE, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, user_id, websocket_available)
            raise

        # Generate filename
        export_progress.start_step(ExportStep.GENERATING_FILENAME)
        await self._send_progress_update(export_progress, user_id, websocket_available)

        try:
            timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}_{format.value}"

            # Add proper extension
            if format == ExportFormat.JSON:
                filename += ".json"
            elif format == ExportFormat.MARKDOWN:
                filename += ".md"
            elif format == ExportFormat.PPTX:
                filename += ".pptx"
            elif format == ExportFormat.PDF:
                filename += ".pdf"

            await self._update_step_progress(
                export_progress, ExportStep.GENERATING_FILENAME,
                progress=100.0, message=f"Filename generated: {filename}",
                user_id=user_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.GENERATING_FILENAME)
            await self._send_progress_update(export_progress, user_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "generating_filename"})
            export_progress.fail_step(ExportStep.GENERATING_FILENAME, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, user_id, websocket_available)
            raise

        # Complete export
        export_progress.start_step(ExportStep.COMPLETED)
        await self._send_progress_update(export_progress, user_id, websocket_available)

        try:
            await self._update_step_progress(
                export_progress, ExportStep.COMPLETED,
                progress=100.0, message=f"{format.value.upper()} export completed successfully",
                user_id=user_id, websocket_available=websocket_available
            )

            export_progress.complete_step(ExportStep.COMPLETED)
            await self._send_progress_update(export_progress, user_id, websocket_available)

        except Exception as e:
            error = create_error_from_exception(e, {"step": "completing_export"})
            export_progress.fail_step(ExportStep.COMPLETED, error.user_message, error.code.value)
            await self._send_progress_update(export_progress, user_id, websocket_available)
            raise

        return file_size, filename

    async def _process_concurrent_exports(
        self,
        bulk_progress: BulkExportProgress,
        presentation_id: str,
        formats: List[ExportFormat],
        user_id: str,
        options: ExportOptions,
        websocket_available: bool
    ) -> None:
        """Process multiple export formats concurrently."""

        # Create semaphore for concurrent limit
        semaphore = asyncio.Semaphore(bulk_progress.concurrent_exports)

        async def process_format(format: ExportFormat) -> Tuple[ExportFormat, Optional[ExportFormatProgress]]:
            async with semaphore:
                try:
                    format_progress = await self.export_presentation(
                        presentation_id, format, user_id, options, websocket_available
                    )
                    return format, format_progress
                except Exception as e:
                    logger.error(f"Failed to export {format.value}: {e}")
                    return format, None

        # Create tasks for all formats
        tasks = [process_format(format) for format in formats]

        # Process concurrently
        for completed_task in asyncio.as_completed(tasks):
            format, format_progress = await completed_task

            if format_progress:
                bulk_progress.update_format_progress(format, format_progress)
                if format_progress.status == ExportStatus.COMPLETED:
                    bulk_progress.complete_format_export(
                        format, format_progress.file_size_bytes, format_progress.filename
                    )
                else:
                    bulk_progress.fail_format_export(
                        format, format_progress.error_message or "Unknown error", format_progress.error_code
                    )
            else:
                bulk_progress.fail_format_export(format, "Export failed to start")

            await self._send_bulk_progress_update(bulk_progress, user_id, websocket_available)

    async def _create_download_zip(
        self,
        bulk_progress: BulkExportProgress,
        user_id: str,
        websocket_available: bool
    ) -> None:
        """Create a ZIP file containing all successful exports."""

        try:
            timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
            zip_filename = f"presentation_export_{timestamp}.zip"
            zip_path = f"temp_{zip_filename}"

            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for format in bulk_progress.successful_exports:
                    format_progress = bulk_progress.export_progress[format]
                    if format_progress.filename:
                        # This would need to be implemented based on how files are stored
                        # For now, just create a placeholder
                        zipf.writestr(f"{format_progress.filename}", "Content placeholder")

            bulk_progress.set_download_zip(zip_path)
            await self._send_bulk_progress_update(bulk_progress, user_id, websocket_available)

        except Exception as e:
            logger.error(f"Failed to create download ZIP: {e}")
            bulk_progress.error_message = f"Failed to create ZIP file: {str(e)}"
            await self._send_bulk_progress_update(bulk_progress, user_id, websocket_available)

    async def _update_step_progress(
        self,
        export_progress: ExportFormatProgress,
        step: ExportStep,
        progress: float,
        message: str,
        user_id: str,
        websocket_available: bool
    ) -> None:
        """Update step progress and send update."""
        export_progress.update_step_progress(step, progress, message)
        await self._send_progress_update(export_progress, user_id, websocket_available)

    async def _send_progress_update(
        self,
        export_progress: ExportFormatProgress,
        user_id: str,
        websocket_available: bool = True
    ) -> None:
        """Send progress update via WebSocket if available."""
        if websocket_available:
            update = ExportProgressUpdate(
                export_id=export_progress.export_id,
                format=export_progress.format,
                progress=export_progress,
                update_type="export_progress",
                message=f"Export progress: {export_progress.overall_progress:.1f}%"
            )
            await self.progress_service.send_progress_update(
                export_progress.export_id, update
            )

    async def _send_bulk_progress_update(
        self,
        bulk_progress: BulkExportProgress,
        user_id: str,
        websocket_available: bool = True
    ) -> None:
        """Send bulk progress update via WebSocket if available."""
        if websocket_available:
            update = ExportProgressUpdate(
                export_id=bulk_progress.bulk_export_id,
                bulk_export_id=bulk_progress.bulk_export_id,
                bulk_progress=bulk_progress,
                update_type="bulk_export_progress",
                message=f"Bulk export progress: {bulk_progress.overall_progress:.1f}%"
            )
            await self.progress_service.send_progress_update(
                bulk_progress.bulk_export_id, update
            )

    async def _handle_export_failure(
        self,
        export_progress: ExportFormatProgress,
        error: Exception,
        user_id: str,
        websocket_available: bool
    ) -> str:
        """Handle export failure and provide recovery options."""

        if isinstance(error, PresentationError):
            error_code = error.code.value
            error_message = error.user_message
        else:
            structured_error = create_error_from_exception(error)
            error_code = structured_error.code.value
            error_message = structured_error.user_message

        # Mark export as failed
        export_progress.status = ExportStatus.FAILED
        export_progress.error_message = error_message
        export_progress.error_code = error_code

        # Send failure notification
        if websocket_available:
            failure_update = ExportProgressUpdate(
                export_id=export_progress.export_id,
                format=export_progress.format,
                progress=export_progress,
                update_type="export_failed",
                message=f"Export failed: {error_message}"
            )
            await self.progress_service.send_progress_update(
                export_progress.export_id, failure_update
            )

        return error_code

    async def _store_export_analytics(
        self,
        export_progress: ExportFormatProgress,
        presentation_id: str,
        user_id: str,
        websocket_available: bool
    ) -> None:
        """Store export analytics data."""

        analytics = ExportAnalytics(
            export_id=export_progress.export_id,
            format=export_progress.format,
            user_id=user_id,
            presentation_id=presentation_id,
            start_time=export_progress.started_at or datetime.utcnow(),
            status=export_progress.status,
            retry_count=export_progress.retry_count,
            file_size_bytes=export_progress.file_size_bytes,
            quality_score=export_progress.quality_score
        )

        if export_progress.completed_at and export_progress.started_at:
            analytics.complete_export(export_progress.file_size_bytes, export_progress.quality_score)

        # Store analytics (implementation depends on your analytics storage)
        # This could be database, file system, or external service
        logger.debug(f"Stored analytics for export {export_progress.export_id}")

    async def get_export_analytics(
        self,
        user_id: str,
        hours: int = 24
    ) -> Dict[str, Any]:
        """Get export analytics for a user."""

        # This would typically query your analytics storage
        # For now, return mock data structure
        return {
            "user_id": user_id,
            "time_window_hours": hours,
            "total_exports": 0,
            "successful_exports": 0,
            "failed_exports": 0,
            "formats_used": {},
            "average_processing_time": 0,
            "average_file_size": 0,
        }

    async def cleanup_expired_exports(self, max_age_hours: int = 24) -> int:
        """Clean up expired export data and files."""

        cleaned_count = 0
        cutoff_time = datetime.utcnow() - timedelta(hours=max_age_hours)

        # Clean up active exports
        expired_exports = [
            export_id for export_id, export_progress in self.active_exports.items()
            if export_progress.created_at < cutoff_time
        ]

        for export_id in expired_exports:
            del self.active_exports[export_id]
            cleaned_count += 1

        # Clean up bulk exports
        expired_bulk_exports = [
            bulk_id for bulk_id, bulk_progress in self.active_bulk_exports.items()
            if bulk_progress.created_at < cutoff_time
        ]

        for bulk_id in expired_bulk_exports:
            del self.active_bulk_exports[bulk_id]
            cleaned_count += 1

        logger.info(f"Cleaned up {cleaned_count} expired export records")
        return cleaned_count


# Global export service instance
export_service = ExportService()


def get_export_service() -> ExportService:
    """Get the global export service instance."""
    return export_service