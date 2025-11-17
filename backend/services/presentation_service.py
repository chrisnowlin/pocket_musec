"""Presentation generation service orchestration.

This service coordinates the end-to-end presentation generation process,
including scaffold building, optional LLM polishing, and persistence.
"""

import json
import logging
from datetime import datetime
from typing import Any, Dict, List, Optional

from backend.repositories.presentation_repository import PresentationRepository
from backend.repositories.lesson_repository import LessonRepository
from backend.llm.chutes_client import ChutesClient
from backend.lessons.presentation_schema import (
    PresentationDocument,
    PresentationStatus,
    PresentationSlide,
    PresentationExport,
)
from backend.lessons.presentation_builder import build_presentation_scaffold
from backend.lessons.presentation_polish import polish_presentation_slides
from backend.lessons.schema_m2 import LessonDocumentM2
from backend.services.presentation_errors import (
    PresentationError,
    PresentationErrorCode,
    PresentationErrorLogger,
    create_error_from_exception,
)
from backend.services.export_status_service import (
    export_status_service,
    ExportFormat,
    ExportStep,
)
from backend.services.style_service import (
    StyleService,
    StyleValidationError,
    StyleNotFoundError,
    StyleAccessDeniedError,
)
from backend.models.style_schema import StyleConfig
from backend.models.preview_schema import PresentationPreview
from pydantic import ValidationError
from typing import TYPE_CHECKING

from backend.utils.casing import camelize

if TYPE_CHECKING:
    from backend.services.preview_service import PreviewService

logger = logging.getLogger(__name__)


class PresentationService:
    """Orchestrates presentation generation from lesson documents."""

    def __init__(
        self,
        presentation_repo: Optional[PresentationRepository] = None,
        lesson_repo: Optional[LessonRepository] = None,
        chutes_client: Optional[ChutesClient] = None,
        style_service: Optional[StyleService] = None,
        preview_service: Optional["PreviewService"] = None,
    ):
        self.presentation_repo = presentation_repo or PresentationRepository()
        self.lesson_repo = lesson_repo or LessonRepository()
        self.chutes_client = chutes_client or ChutesClient(require_api_key=False)
        self.style_service = style_service or StyleService()
        # Don't instantiate PreviewService by default to avoid circular dependency
        self.preview_service = preview_service

    def generate_presentation(
        self,
        lesson_id: str,
        user_id: str,
        style: str = "default",
        use_llm_polish: bool = True,
        timeout_seconds: int = 30,
    ) -> PresentationDocument:
        """Generate a presentation for a lesson.

        Args:
            lesson_id: ID of the lesson to generate presentation for
            user_id: ID of the user requesting the presentation
            style: Presentation style (currently only 'default' supported)
            use_llm_polish: Whether to attempt LLM polishing
            timeout_seconds: Maximum time to wait for LLM polishing

        Returns:
            PresentationDocument with generated slides

        Raises:
            PresentationError: If presentation generation fails
        """
        logger.info(f"Starting presentation generation for lesson {lesson_id}")

        try:
            # Validate inputs
            self._validate_generation_params(lesson_id, user_id, style, timeout_seconds)

            # 1. Fetch the lesson document
            lesson = self._fetch_lesson_document_with_error_handling(lesson_id, user_id)

            # 2. Mark any existing presentations as stale
            self._mark_existing_stale(lesson_id, lesson.revision)

            # 3. Generate scaffold slides
            try:
                scaffold_slides = build_presentation_scaffold(lesson)
            except Exception as e:
                error = PresentationError.internal_error(
                    f"Scaffold generation failed: {str(e)}"
                )
                PresentationErrorLogger.log_error(error, {"lesson_id": lesson_id})
                raise error

            # 4. Check if presentation already exists for this revision and update it
            presentation = self._create_or_update_presentation(
                lesson_id, lesson.revision, scaffold_slides, style, user_id
            )

            # 5. Generate slides (scaffold + optional polish)
            final_slides = self._generate_slides_with_error_handling(
                lesson, scaffold_slides, use_llm_polish, timeout_seconds
            )

            # 6. Update presentation with final slides
            updated_presentation = self._update_presentation_with_error_handling(
                presentation.id, final_slides
            )

            # 7. Generate initial export assets
            self._generate_export_assets_with_error_handling(updated_presentation)

            logger.info(f"Successfully generated presentation {presentation.id}")
            return updated_presentation

        except PresentationError as e:
            # Log and re-raise structured errors
            PresentationErrorLogger.log_error(
                e,
                {
                    "lesson_id": lesson_id,
                    "user_id": user_id,
                    "style": style,
                    "use_llm_polish": use_llm_polish,
                    "timeout_seconds": timeout_seconds,
                },
            )

            # Mark presentation as error if it exists
            try:
                if "presentation" in locals():
                    self.presentation_repo.update_presentation_status(
                        presentation_id=presentation.id,
                        status=PresentationStatus.ERROR,
                        error_code=e.code.value,
                        error_message=e.user_message,
                    )
            except Exception as db_error:
                logger.error(f"Failed to mark presentation as error: {db_error}")

            raise

        except Exception as e:
            # Convert unstructured exceptions to structured errors
            structured_error = create_error_from_exception(
                e,
                {
                    "lesson_id": lesson_id,
                    "user_id": user_id,
                    "operation": "generate_presentation",
                },
            )

            PresentationErrorLogger.log_error(structured_error)

            # Mark presentation as error if it exists
            try:
                if "presentation" in locals():
                    self.presentation_repo.update_presentation_status(
                        presentation_id=presentation.id,
                        status=PresentationStatus.ERROR,
                        error_code=structured_error.code.value,
                        error_message=structured_error.user_message,
                    )
            except Exception as db_error:
                logger.error(f"Failed to mark presentation as error: {db_error}")

            raise structured_error

    def regenerate_presentation(
        self,
        lesson_id: str,
        user_id: str,
        style: str = "default",
        use_llm_polish: bool = True,
        timeout_seconds: int = 30,
    ) -> PresentationDocument:
        """Regenerate a presentation for a lesson.

        Similar to generate_presentation but explicitly marks all existing
        presentations as stale before creating a new one.
        """
        logger.info(f"Regenerating presentation for lesson {lesson_id}")

        # Validate inputs
        self._validate_generation_params(lesson_id, user_id, style, timeout_seconds)

        # Mark all existing presentations as stale
        self._mark_all_stale(lesson_id)

        # Generate new presentation
        return self.generate_presentation(
            lesson_id, user_id, style, use_llm_polish, timeout_seconds
        )

    def get_presentation_status(self, lesson_id: str) -> Optional[Dict[str, Any]]:
        """Return the latest presentation for a lesson as a status dict."""
        if not lesson_id or not isinstance(lesson_id, str):
            raise PresentationError.validation_failed("lesson_id", lesson_id)

        presentation = self.presentation_repo.latest_by_lesson(lesson_id)
        if not presentation:
            return None

        return self.serialize_presentation_summary(presentation)

    def serialize_presentation_summary(
        self, presentation: PresentationDocument
    ) -> Dict[str, Any]:
        """Convert a PresentationDocument into a camelCase summary payload."""

        metadata = self._build_presentation_metadata(presentation)

        status_payload = {
            "id": presentation.id,
            "presentation_id": presentation.id,  # backwards compatibility for UI
            "lesson_id": presentation.lesson_id,
            "lesson_revision": presentation.lesson_revision,
            "version": presentation.version,
            "status": presentation.status.value,
            "style": presentation.style,
            "slide_count": len(presentation.slides),
            "created_at": presentation.created_at.isoformat(),
            "updated_at": presentation.updated_at.isoformat(),
            "has_exports": len(presentation.export_assets) > 0,
            "error_code": presentation.error_code,
            "error_message": presentation.error_message,
            **metadata,
        }
        return camelize(status_payload)

    def get_presentation(
        self, presentation_id: str, user_id: str
    ) -> Optional[PresentationDocument]:
        """Get a presentation by ID with user access check."""
        if not presentation_id or not isinstance(presentation_id, str):
            raise PresentationError.validation_failed(
                "presentation_id", presentation_id
            )
        if not user_id or not isinstance(user_id, str):
            raise PresentationError.validation_failed("user_id", user_id)

        presentation = self.presentation_repo.get_presentation(presentation_id)
        if not presentation:
            return None

        lesson = self.lesson_repo.get_lesson(presentation.lesson_id)
        if not lesson or lesson.user_id != user_id:
            return None

        return presentation

    def mark_stale_on_lesson_update(self, lesson_id: str, new_revision: int) -> int:
        """Mark presentations as stale when a lesson is updated.

        Args:
            lesson_id: ID of the updated lesson
            new_revision: New revision number of the lesson

        Returns:
            Number of presentations marked as stale
        """
        if not lesson_id or not isinstance(lesson_id, str):
            raise PresentationError.validation_failed("lesson_id", lesson_id)
        if not isinstance(new_revision, int) or new_revision < 0:
            raise PresentationError.validation_failed("new_revision", new_revision)

        return self.presentation_repo.mark_stale_for_lesson(lesson_id, new_revision)

    def generate_preview(
        self,
        presentation_id: str,
        style_id: Optional[str] = None,
    ) -> PresentationPreview:
        """Generate a lightweight preview for a presentation.

        Args:
            presentation_id: ID of the presentation to preview.
            style_id: Optional style identifier to apply for the preview.

        Returns:
            PresentationPreview model containing slide outlines and estimated duration.

        Raises:
            PresentationError: If preview generation fails or presentation not found.
        """
        try:
            # Retrieve existing presentation metadata
            presentation = self.presentation_repo.get_presentation(presentation_id)
            if not presentation:
                raise PresentationError.presentation_not_found(presentation_id)

            # Use the preview service to generate a preview based on current data
            preview = self.preview_service.generate_preview(
                presentation_id=presentation_id,
                style_id=style_id,
            )
            return preview
        except PresentationError:
            raise
        except Exception as e:
            # Convert unexpected errors to structured presentation errors
            structured_error = create_error_from_exception(
                e,
                {
                    "presentation_id": presentation_id,
                    "operation": "generate_preview",
                },
            )
            raise structured_error

        """Get and validate a style configuration.

        Args:
            style_id_or_name: Style ID or name
            user_id: Optional user ID for access control

        Returns:
            Validated StyleConfig

        Raises:
            PresentationError: If style validation fails
        """
        try:
            return self.style_service.validate_and_apply_style(
                style_id_or_name, user_id
            )
        except (StyleValidationError, StyleNotFoundError, StyleAccessDeniedError) as e:
            if isinstance(e, StyleValidationError):
                raise PresentationError.invalid_style(str(e))
            elif isinstance(e, StyleNotFoundError):
                raise PresentationError.style_not_found(style_id_or_name)
            elif isinstance(e, StyleAccessDeniedError):
                raise PresentationError.style_access_denied(style_id_or_name)
        except Exception as e:
            raise PresentationError.internal_error(f"Error getting style: {str(e)}")

    def _build_presentation_metadata(
        self, presentation: PresentationDocument
    ) -> Dict[str, Any]:
        """Compute derived presentation metadata for API responses."""
        title = None
        description = None

        if presentation.slides:
            first_slide = presentation.slides[0]
            title = first_slide.title

            if first_slide.teacher_script:
                description = first_slide.teacher_script[:280]
            elif first_slide.key_points:
                description = "; ".join(first_slide.key_points)[:280]

        total_duration = sum(
            slide.duration_minutes or 0 for slide in presentation.slides
        )

        return {
            "title": title,
            "description": description,
            "total_slides": len(presentation.slides),
            "total_duration_minutes": total_duration or None,
            "is_stale": presentation.status == PresentationStatus.STALE,
        }

    def _validate_generation_params(
        self, lesson_id: str, user_id: str, style: str, timeout_seconds: int
    ) -> None:
        """Validate presentation generation parameters."""
        if not lesson_id or not isinstance(lesson_id, str):
            raise PresentationError.validation_failed("lesson_id", lesson_id)

        if not user_id or not isinstance(user_id, str):
            raise PresentationError.validation_failed("user_id", user_id)

        # Validate style using style service
        try:
            # Handle both string IDs and StyleConfig objects
            if isinstance(style, str):
                self.style_service.validate_and_apply_style(style, user_id)
            else:
                self.style_service.validate_and_apply_style(style, user_id)
        except StyleNotFoundError:
            # Check if it's a valid preset ID
            try:
                self.style_service.get_preset_style(style)
            except StyleNotFoundError:
                raise PresentationError.validation_failed("style", style)
        except StyleValidationError as e:
            raise PresentationError.validation_failed(
                "style", f"Invalid style: {str(e)}"
            )
        except StyleAccessDeniedError as e:
            raise PresentationError.permission_denied(f"No access to style: {str(e)}")

        if (
            not isinstance(timeout_seconds, int)
            or timeout_seconds < 1
            or timeout_seconds > 300
        ):
            raise PresentationError.validation_failed(
                "timeout_seconds", timeout_seconds
            )

    def _fetch_lesson_document_with_error_handling(
        self, lesson_id: str, user_id: str
    ) -> LessonDocumentM2:
        """Fetch lesson document with proper error handling."""
        try:
            lesson = self.lesson_repo.get_lesson(lesson_id)
            if not lesson:
                raise PresentationError.lesson_not_found(lesson_id)

            if lesson.user_id != user_id:
                raise PresentationError.lesson_access_denied(lesson_id)

            document = self._fetch_lesson_document(lesson_id, user_id)
            if not document:
                # Try to create fallback document
                try:
                    document = self._create_fallback_lesson_document(lesson)
                except Exception as e:
                    error = PresentationError.lesson_parse_failed(
                        f"Failed to create fallback document: {str(e)}"
                    )
                    error.context = {"lesson_id": lesson_id, "user_id": user_id}
                    raise error

            return document

        except PresentationError:
            raise
        except Exception as e:
            error = create_error_from_exception(
                e,
                {
                    "lesson_id": lesson_id,
                    "user_id": user_id,
                    "operation": "fetch_lesson_document",
                },
            )
            raise error

    def _create_or_update_presentation(
        self,
        lesson_id: str,
        lesson_revision: int,
        scaffold_slides: List[PresentationSlide],
        style: str,
        user_id: str,
    ) -> PresentationDocument:
        """Create or update presentation with error handling."""
        try:
            existing_presentation = self.presentation_repo.latest_by_lesson(lesson_id)

            if (
                existing_presentation
                and existing_presentation.lesson_revision == lesson_revision
            ):
                # Update existing presentation
                presentation = self.presentation_repo.update_presentation_slides(
                    presentation_id=existing_presentation.id,
                    slides=scaffold_slides,
                )
                # Reset status to pending
                self.presentation_repo.update_presentation_status(
                    presentation_id=presentation.id,
                    status=PresentationStatus.PENDING,
                    error_code=None,
                    error_message=None,
                )
            else:
                # Create new presentation
                presentation = self.presentation_repo.create_presentation(
                    lesson_id=lesson_id,
                    lesson_revision=lesson_revision,
                    slides=scaffold_slides,
                    style=style,
                )

            return presentation

        except PresentationError:
            raise
        except Exception as e:
            error = create_error_from_exception(
                e,
                {
                    "lesson_id": lesson_id,
                    "user_id": user_id,
                    "operation": "create_or_update_presentation",
                },
            )
            raise error

    def _generate_slides_with_error_handling(
        self,
        lesson: LessonDocumentM2,
        scaffold_slides: List[PresentationSlide],
        use_llm_polish: bool,
        timeout_seconds: int,
    ) -> List[PresentationSlide]:
        """Generate slides with comprehensive error handling."""
        if not use_llm_polish or not self.chutes_client.is_available():
            logger.info("Skipping LLM polishing - not available or disabled")
            return scaffold_slides

        try:
            return polish_presentation_slides(
                lesson=lesson,
                scaffold_slides=scaffold_slides,
                chutes_client=self.chutes_client,
                timeout_seconds=timeout_seconds,
            )

        except TimeoutError as e:
            error = PresentationError.llm_timeout(timeout_seconds)
            PresentationErrorLogger.log_error(
                error,
                {
                    "lesson_id": lesson.id,
                    "timeout_seconds": timeout_seconds,
                },
            )
            # Return scaffold slides as fallback
            PresentationErrorLogger.log_recovery_attempt(error, "using scaffold slides")
            return scaffold_slides

        except ConnectionError as e:
            error = PresentationError.llm_unavailable()
            PresentationErrorLogger.log_error(
                error,
                {
                    "lesson_id": lesson.id,
                    "error_details": str(e),
                },
            )
            # Return scaffold slides as fallback
            PresentationErrorLogger.log_recovery_attempt(error, "using scaffold slides")
            return scaffold_slides

        except Exception as e:
            # Check for rate limiting
            if "rate" in str(e).lower() or "limit" in str(e).lower():
                error = PresentationError.llm_rate_limited()
                PresentationErrorLogger.log_error(
                    error,
                    {
                        "lesson_id": lesson.id,
                        "error_details": str(e),
                    },
                )
                # Retry once after a delay if available
                PresentationErrorLogger.log_recovery_attempt(error, "retrying once")
                try:
                    import time

                    time.sleep(2)  # Brief delay
                    return polish_presentation_slides(
                        lesson=lesson,
                        scaffold_slides=scaffold_slides,
                        chutes_client=self.chutes_client,
                        timeout_seconds=timeout_seconds,
                    )
                except Exception:
                    # If retry fails, return scaffold slides
                    return scaffold_slides
            else:
                # Log and return scaffold as fallback
                error = PresentationError.llm_invalid_response()
                error.technical_message = str(e)
                PresentationErrorLogger.log_error(
                    error,
                    {
                        "lesson_id": lesson.id,
                        "error_details": str(e),
                    },
                )
                PresentationErrorLogger.log_recovery_attempt(
                    error, "using scaffold slides"
                )
                return scaffold_slides

    def _update_presentation_with_error_handling(
        self, presentation_id: str, slides: List[PresentationSlide]
    ) -> PresentationDocument:
        """Update presentation with slides and error handling."""
        try:
            return self.presentation_repo.update_presentation_slides(
                presentation_id=presentation_id,
                slides=slides,
            )

        except PresentationError:
            raise
        except Exception as e:
            error = create_error_from_exception(
                e,
                {
                    "presentation_id": presentation_id,
                    "operation": "update_presentation_slides",
                },
            )
            raise error

    def _generate_export_assets_with_error_handling(
        self, presentation: PresentationDocument
    ) -> None:
        """Generate export assets with comprehensive error handling and progress tracking."""
        try:
            self._generate_export_assets_with_progress(presentation)
        except Exception as e:
            # Don't fail the whole presentation generation for export issues
            error = create_error_from_exception(
                e,
                {
                    "presentation_id": presentation.id,
                    "operation": "generate_export_assets",
                },
            )
            PresentationErrorLogger.log_error(error)
            PresentationErrorLogger.log_recovery_attempt(
                error, "skipping export generation"
            )

    def _generate_export_assets_with_progress(
        self, presentation: PresentationDocument
    ) -> None:
        """Generate export assets with progress tracking."""
        logger.info(
            f"Starting export generation with progress tracking for presentation {presentation.id}"
        )

        # Generate each export format with individual job tracking
        export_formats = [
            (ExportFormat.JSON, self._create_json_export, "JSON"),
            (ExportFormat.MARKDOWN, self._create_markdown_export, "Markdown"),
            (ExportFormat.PPTX, self._create_pptx_export, "PowerPoint"),
            (ExportFormat.PDF, self._create_pdf_export, "PDF"),
        ]

        for format_enum, creator_func, format_name in export_formats:
            try:
                logger.info(f"Starting {format_name} export generation")

                # Queue export with progress tracking
                job_id = export_status_service.queue_export_job(
                    presentation_id=presentation.id,
                    export_format=format_enum,
                    executor_callback=lambda f=creator_func, p=presentation: f(p),
                    priority=1,
                )

                logger.info(f"Queued {format_name} export with job_id: {job_id}")

            except Exception as e:
                logger.error(f"Failed to queue {format_name} export: {e}")
                # Continue with other formats even if one fails
                continue

    def generate_export_with_status(
        self, presentation_id: str, export_format: str
    ) -> str:
        """Generate a single export format and return the job ID for status tracking."""
        presentation = self.presentation_repo.get_presentation(presentation_id)
        if not presentation:
            raise PresentationError.validation_failed(
                "presentation_id", presentation_id
            )

        # Parse export format
        try:
            format_enum = ExportFormat(export_format.lower())
        except ValueError:
            raise PresentationError.validation_failed("export_format", export_format)

        # Create export callback based on format
        export_callbacks = {
            ExportFormat.JSON: self._create_json_export,
            ExportFormat.MARKDOWN: self._create_markdown_export,
            ExportFormat.PPTX: self._create_pptx_export,
            ExportFormat.PDF: self._create_pdf_export,
        }

        callback = export_callbacks.get(format_enum)
        if not callback:
            raise PresentationError.validation_failed("export_format", export_format)

        # Queue export with high priority
        job_id = export_status_service.queue_export_job(
            presentation_id=presentation_id,
            export_format=format_enum,
            executor_callback=lambda: callback(presentation),
            priority=5,  # High priority for on-demand exports
        )

        logger.info(f"Started on-demand export job {job_id} for format {export_format}")
        return job_id

    def get_export_status(
        self, presentation_id: str, export_format: Optional[str] = None
    ) -> Dict[str, Any]:
        """Get export status information for a presentation."""
        jobs = export_status_service.get_jobs_for_presentation(presentation_id)

        # Filter by format if specified
        if export_format:
            try:
                format_enum = ExportFormat(export_format.lower())
                jobs = [job for job in jobs if job.export_format == format_enum]
            except ValueError:
                pass

        # Calculate overall status
        result = {
            "presentation_id": presentation_id,
            "exports": [],
            "overall_status": "not_started",
            "total_jobs": len(jobs),
            "completed_jobs": len([j for j in jobs if j.status == "completed"]),
            "failed_jobs": len([j for j in jobs if j.status == "failed"]),
            "active_jobs": len(
                [j for j in jobs if j.status in ["pending", "processing"]]
            ),
        }

        for job in jobs:
            export_info = {
                "job_id": job.job_id,
                "format": job.export_format.value,
                "status": job.status.value,
                "progress": {
                    "percentage": job.progress.progress_percentage,
                    "current_step": job.progress.current_step.value,
                    "estimated_time_remaining": job.progress.estimated_time_remaining,
                },
                "created_at": job.created_at.isoformat(),
                "started_at": job.started_at.isoformat() if job.started_at else None,
                "completed_at": job.completed_at.isoformat()
                if job.completed_at
                else None,
                "error_message": job.error_message,
                "retry_count": job.retry_count,
                "export_asset": {
                    "url_or_path": job.export_asset.url_or_path,
                    "file_size_bytes": job.export_asset.file_size_bytes,
                    "generated_at": job.export_asset.generated_at.isoformat()
                    if job.export_asset
                    else None,
                }
                if job.export_asset
                else None,
            }
            result["exports"].append(export_info)

        # Determine overall status
        if (
            result["completed_jobs"] == result["total_jobs"]
            and result["total_jobs"] > 0
        ):
            result["overall_status"] = "completed"
        elif result["failed_jobs"] == result["total_jobs"] and result["total_jobs"] > 0:
            result["overall_status"] = "failed"
        elif result["active_jobs"] > 0:
            result["overall_status"] = "processing"
        elif result["total_jobs"] == 0:
            result["overall_status"] = "not_started"
        else:
            result["overall_status"] = "partial"

        return result

    def retry_failed_export(
        self, presentation_id: str, export_format: str
    ) -> Optional[str]:
        """Retry a failed export and return the new job ID."""
        format_enum = ExportFormat(export_format.lower())
        jobs = export_status_service.get_jobs_for_presentation(presentation_id)

        # Find the most recent failed job for this format
        failed_jobs = [
            job
            for job in jobs
            if job.export_format == format_enum and job.status == "failed"
        ]

        if not failed_jobs:
            return None

        # Sort by creation time to get the most recent
        failed_jobs.sort(key=lambda j: j.created_at, reverse=True)
        last_job = failed_jobs[0]

        logger.info(f"Retrying failed export job {last_job.job_id}")

        # Create new export job
        return self.generate_export_with_status(presentation_id, export_format)

    def _fetch_lesson_document(
        self, lesson_id: str, user_id: str
    ) -> Optional[LessonDocumentM2]:
        """Fetch and validate lesson document for the user."""
        lesson = self.lesson_repo.get_lesson(lesson_id)
        if not lesson or lesson.user_id != user_id:
            return None

        metadata_doc: Optional[LessonDocumentM2] = None
        if lesson.metadata:
            try:
                metadata = json.loads(lesson.metadata)
                raw_doc = metadata.get("lesson_document")
                if raw_doc:
                    validate_fn = getattr(LessonDocumentM2, "model_validate", None)
                    metadata_doc = (
                        validate_fn(raw_doc)
                        if callable(validate_fn)
                        else LessonDocumentM2(**raw_doc)
                    )
            except (json.JSONDecodeError, ValidationError) as exc:
                logger.warning(
                    "Failed to parse lesson %s metadata as m2.0: %s", lesson_id, exc
                )

        if metadata_doc:
            return metadata_doc

        # Fallback 1: attempt to treat lesson.content as JSON for edge cases
        try:
            lesson_data = json.loads(lesson.content)
            return LessonDocumentM2(**lesson_data)
        except Exception as exc:
            logger.debug("Lesson %s content is not JSON: %s", lesson_id, exc)

        # Fallback 2: convert markdown content to minimal structured document
        try:
            return self._create_fallback_lesson_document(lesson)
        except Exception as exc:
            logger.warning("Lesson %s lacks structured document: %s", lesson_id, exc)
            return None

    def _create_fallback_lesson_document(self, lesson) -> LessonDocumentM2:
        """Create a minimal structured lesson document from markdown content."""
        import re
        import json
        from datetime import datetime, timezone

        # Parse basic information from lesson
        title = lesson.title or "Untitled Lesson"

        # Extract grade and strand from metadata JSON if available
        grade = "All Grades"
        strand = "All Strands"
        if lesson.metadata:
            try:
                metadata = json.loads(lesson.metadata)
                grade = metadata.get("grade_level", "All Grades")
                strand = metadata.get("strand_code", "All Strands")
            except (json.JSONDecodeError, AttributeError):
                pass  # Use defaults

        # Extract objectives from markdown content
        objectives = []
        content_lines = lesson.content.split("\n")

        # Look for learning objectives section
        in_objectives = False
        for line in content_lines:
            line = line.strip()
            if "learning objective" in line.lower() or "objective" in line.lower():
                in_objectives = True
                continue
            elif line.startswith("##") and in_objectives:
                # Hit next section, stop collecting objectives
                break
            elif in_objectives and line and not line.startswith("#"):
                # Clean up objective text
                obj = re.sub(r"^[\-\*\d\.\s]+", "", line).strip()
                if obj and len(obj) > 10:
                    objectives.append(obj)

        # If no objectives found, create a default one
        if not objectives:
            objectives = [f"Students will engage with the lesson content: {title}"]

        # Extract materials if present
        materials = []
        in_materials = False
        for line in content_lines:
            line = line.strip()
            if "material" in line.lower() or "resource" in line.lower():
                in_materials = True
                continue
            elif line.startswith("##") and in_materials:
                break
            elif in_materials and line and not line.startswith("#"):
                material = re.sub(r"^[\-\*\d\.\s]+", "", line).strip()
                if material and len(material) > 5:
                    materials.append(material)

        # Create timing (use default 45 minutes)
        timing = {"total_minutes": 45}

        # Create content structure
        content = {
            "materials": materials,
            "warmup": "Introduction and lesson overview",
            "activities": [
                {
                    "id": "main_activity",
                    "title": "Main Learning Activity",
                    "duration_minutes": 30,
                    "steps": [
                        "Introduction",
                        "Guided Practice",
                        "Independent Work",
                        "Reflection",
                    ],
                    "aligned_standards": [],
                    "citations": [],
                }
            ],
            "assessment": "Formative assessment through observation and questioning",
            "differentiation": "Differentiated instruction as needed for diverse learners",
            "exit_ticket": "Lesson reflection and key takeaways",
            "timing": timing,
        }

        # Create minimal lesson document
        # Parse datetime - handle the format from database (ISO string without timezone)
        try:
            # Ensure we have string values
            created_at_str = str(lesson.created_at) if lesson.created_at else None
            updated_at_str = str(lesson.updated_at) if lesson.updated_at else None

            if created_at_str:
                created_at = datetime.fromisoformat(created_at_str)
            else:
                created_at = datetime.now(timezone.utc)

            if updated_at_str:
                updated_at = datetime.fromisoformat(updated_at_str)
            else:
                updated_at = datetime.now(timezone.utc)

        except (ValueError, AttributeError, TypeError):
            # Fallback to current time if parsing fails
            created_at = datetime.now(timezone.utc)
            updated_at = datetime.now(timezone.utc)

        lesson_doc = {
            "id": lesson.id,
            "created_at": created_at,
            "updated_at": updated_at,
            "version": "m2.0",
            "title": title,
            "grade": grade,
            "strands": [strand],
            "standards": [],
            "objectives": objectives[:5],  # Limit to 5 objectives
            "content": content,
            "citations": [],
            "revision": 1,
        }

        return LessonDocumentM2(**lesson_doc)

    def _mark_existing_stale(self, lesson_id: str, current_revision: int) -> None:
        """Mark existing presentations as stale for older revisions."""
        stale_count = self.presentation_repo.mark_stale_for_lesson(
            lesson_id, current_revision
        )
        if stale_count > 0:
            logger.info(f"Marked {stale_count} existing presentations as stale")

    def _mark_all_stale(self, lesson_id: str) -> None:
        """Mark all presentations for a lesson as stale."""
        # This is a more aggressive stale marking for regeneration
        presentations = self.presentation_repo.list_presentations_for_lesson(
            lesson_id, include_stale=False
        )

        for presentation in presentations:
            self.presentation_repo.update_presentation_status(
                presentation_id=presentation.id,
                status=PresentationStatus.STALE,
            )

        logger.info(
            f"Marked {len(presentations)} presentations as stale for regeneration"
        )

    def _generate_slides(
        self,
        lesson: LessonDocumentM2,
        scaffold_slides: List,
        use_llm_polish: bool,
        timeout_seconds: int,
    ) -> List:
        """Generate final slides from lesson content."""
        if use_llm_polish and self.chutes_client.is_available():
            try:
                return polish_presentation_slides(
                    lesson=lesson,
                    scaffold_slides=scaffold_slides,
                    chutes_client=self.chutes_client,
                    timeout_seconds=timeout_seconds,
                )
            except Exception as e:
                logger.warning(f"LLM polishing failed, using scaffold: {e}")

        return scaffold_slides

    def _generate_export_assets(self, presentation: PresentationDocument) -> None:
        """Generate initial export assets for the presentation."""
        try:
            # Generate JSON export
            json_export = self._create_json_export(presentation)
            self.presentation_repo.add_export_asset(
                presentation_id=presentation.id,
                export_asset=json_export,
            )

            # Generate Markdown export
            markdown_export = self._create_markdown_export(presentation)
            self.presentation_repo.add_export_asset(
                presentation_id=presentation.id,
                export_asset=markdown_export,
            )

            # Generate PPTX export
            try:
                pptx_export = self._create_pptx_export(presentation)
                self.presentation_repo.add_export_asset(
                    presentation_id=presentation.id,
                    export_asset=pptx_export,
                )
            except Exception as pptx_error:
                logger.warning(f"PPTX export failed: {pptx_error}")

            # Generate PDF export
            try:
                pdf_export = self._create_pdf_export(presentation)
                self.presentation_repo.add_export_asset(
                    presentation_id=presentation.id,
                    export_asset=pdf_export,
                )
            except Exception as pdf_error:
                logger.warning(f"PDF export failed: {pdf_error}")

            logger.info(f"Generated export assets for presentation {presentation.id}")

        except Exception as e:
            logger.warning(f"Failed to generate export assets: {e}")
            # Don't fail the whole presentation generation for export issues

    def _create_json_export(
        self, presentation: PresentationDocument
    ) -> PresentationExport:
        """Create JSON export asset."""
        import json

        # Use a custom encoder to handle datetime objects
        class DateTimeEncoder(json.JSONEncoder):
            def default(self, obj):
                if isinstance(obj, datetime):
                    return obj.isoformat()
                return super().default(obj)

        json_content = json.dumps(
            {
                "presentation": presentation.model_dump(mode="json"),
                "generated_at": datetime.utcnow().isoformat(),
                "format": "p1.0",
            },
            indent=2,
            cls=DateTimeEncoder,
        )

        # Create export asset with proper datetime handling
        export_datetime = datetime.utcnow()
        return PresentationExport(
            format="json",
            url_or_path=f"presentation_{presentation.id}.json",
            generated_at=export_datetime,
            file_size_bytes=len(json_content.encode("utf-8")),
        )

    def _create_markdown_export(
        self, presentation: PresentationDocument
    ) -> PresentationExport:
        """Create Markdown export asset."""
        markdown_lines = [
            f"# {presentation.slides[0].title if presentation.slides else 'Presentation'}",
            "",
            f"**Generated:** {presentation.created_at.strftime('%Y-%m-%d %H:%M')}",
            f"**Lesson ID:** {presentation.lesson_id}",
            f"**Revision:** {presentation.lesson_revision}",
            "",
            "---",
            "",
        ]

        for slide in presentation.slides:
            markdown_lines.extend(
                [
                    f"## {slide.title}",
                    "",
                ]
            )

            if slide.subtitle:
                markdown_lines.extend(
                    [
                        f"**{slide.subtitle}**",
                        "",
                    ]
                )

            if slide.key_points:
                markdown_lines.extend(
                    [
                        "**Key Points:**",
                        "",
                    ]
                )
                for point in slide.key_points:
                    markdown_lines.append(f"- {point}")
                markdown_lines.append("")

            if slide.teacher_script:
                markdown_lines.extend(
                    [
                        "**Teacher Script:**",
                        "",
                        slide.teacher_script,
                        "",
                    ]
                )

            if slide.duration_minutes:
                markdown_lines.extend(
                    [
                        f"**Duration:** {slide.duration_minutes} minutes",
                        "",
                    ]
                )

            if slide.standard_codes:
                markdown_lines.extend(
                    [
                        "**Standards:** " + ", ".join(slide.standard_codes),
                        "",
                    ]
                )

            markdown_lines.append("---")
            markdown_lines.append("")

        markdown_content = "\n".join(markdown_lines)

        # Create export asset with proper datetime handling
        export_datetime = datetime.utcnow()
        return PresentationExport(
            format="markdown",
            url_or_path=f"presentation_{presentation.id}.md",
            generated_at=export_datetime,
            file_size_bytes=len(markdown_content.encode("utf-8")),
        )

    def _create_pptx_export(
        self, presentation: PresentationDocument
    ) -> PresentationExport:
        """Create PowerPoint PPTX export asset."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            import tempfile
            import os
        except ModuleNotFoundError as e:
            error = PresentationError.export_failed(
                "pptx", f"Missing dependency: {e.name}"
            )
            raise error
        except ImportError as e:
            error = PresentationError.export_failed("pptx", f"Import error: {str(e)}")
            raise error

        try:
            # Create presentation
            prs = Presentation()

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = (
                presentation.slides[0].title if presentation.slides else "Presentation"
            )
            subtitle.text = (
                f"Generated: {presentation.created_at.strftime('%Y-%m-%d %H:%M')}"
            )

            # Add content slides
            for content_slide in presentation.slides[
                1:
            ]:  # Skip title slide if it's the first one
                # Choose layout based on content
                if content_slide.title and content_slide.key_points:
                    slide_layout = prs.slide_layouts[1]  # Title and content
                else:
                    slide_layout = prs.slide_layouts[5]  # Blank

                slide = prs.slides.add_slide(slide_layout)

                # Add title
                if content_slide.title and slide.shapes.title:
                    slide.shapes.title.text = content_slide.title

                # Add content
                content_parts = []

                if content_slide.key_points:
                    content_parts.extend(
                        [f"• {point}" for point in content_slide.key_points]
                    )

                if content_slide.teacher_script:
                    content_parts.append(f"Teacher: {content_slide.teacher_script}")

                # Add content to slide if we have a content placeholder
                if content_parts and len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    content_placeholder.text = "\n".join(content_parts)

            # Save to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
                prs.save(temp_file.name)

                # Get file size
                file_size = os.path.getsize(temp_file.name)

                # Move to permanent location
                final_path = f"presentation_{presentation.id}.pptx"
                os.rename(temp_file.name, final_path)

                # Create export asset with proper datetime handling
                export_datetime = datetime.utcnow()
                return PresentationExport(
                    format="pptx",
                    url_or_path=final_path,
                    generated_at=export_datetime,
                    file_size_bytes=file_size,
                )

        except PermissionError as e:
            error = PresentationError.export_permission_denied()
            error.technical_message = f"File permission error: {str(e)}"
            raise error
        except OSError as e:
            error = PresentationError.export_storage_failed()
            error.technical_message = f"File system error: {str(e)}"
            raise error
        except Exception as e:
            error = PresentationError.export_failed("pptx", str(e))
            raise error

    def _create_pdf_export(
        self, presentation: PresentationDocument
    ) -> PresentationExport:
        """Create PDF export asset."""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.units import inch
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.colors import black, blue
            import tempfile
            import os
        except ModuleNotFoundError as e:
            error = PresentationError.export_failed(
                "pdf", f"Missing dependency: {e.name}"
            )
            raise error
        except ImportError as e:
            error = PresentationError.export_failed("pdf", f"Import error: {str(e)}")
            raise error

        try:
            # Create temporary PDF file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
                temp_filename = temp_file.name

            # Create PDF document
            doc = SimpleDocTemplate(temp_filename, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []

            # Add title
            title_style = styles["Title"]
            title_style.textColor = blue
            title = (
                presentation.slides[0].title if presentation.slides else "Presentation"
            )
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 12))

            # Add metadata
            normal_style = styles["Normal"]
            story.append(
                Paragraph(
                    f"<b>Generated:</b> {presentation.created_at.strftime('%Y-%m-%d %H:%M')}",
                    normal_style,
                )
            )
            story.append(
                Paragraph(
                    f"<b>Total Slides:</b> {len(presentation.slides)}", normal_style
                )
            )
            story.append(Spacer(1, 20))

            # Add slides content
            for i, slide in enumerate(presentation.slides):
                # Slide title
                if slide.title:
                    heading_style = styles["Heading1"]
                    story.append(
                        Paragraph(f"Slide {i + 1}: {slide.title}", heading_style)
                    )
                    story.append(Spacer(1, 12))

                # Key points
                if slide.key_points:
                    story.append(Paragraph("<b>Key Points:</b>", normal_style))
                    for point in slide.key_points:
                        story.append(Paragraph(f"• {point}", normal_style))
                    story.append(Spacer(1, 6))

                # Teacher script
                if slide.teacher_script:
                    story.append(Paragraph("<b>Teacher Script:</b>", normal_style))
                    story.append(Paragraph(slide.teacher_script, normal_style))
                    story.append(Spacer(1, 6))

                # Duration
                if slide.duration_minutes:
                    story.append(
                        Paragraph(
                            f"<b>Duration:</b> {slide.duration_minutes} minutes",
                            normal_style,
                        )
                    )
                    story.append(Spacer(1, 6))

                # Standards
                if slide.standard_codes:
                    story.append(
                        Paragraph(
                            f"<b>Standards:</b> {', '.join(slide.standard_codes)}",
                            normal_style,
                        )
                    )
                    story.append(Spacer(1, 6))

                # Add separator between slides
                story.append(Spacer(1, 20))

            # Build PDF
            doc.build(story)

            # Get file size
            file_size = os.path.getsize(temp_filename)

            # Move to permanent location
            final_path = f"presentation_{presentation.id}.pdf"
            os.rename(temp_filename, final_path)

            # Create export asset with proper datetime handling
            export_datetime = datetime.utcnow()
            return PresentationExport(
                format="pdf",
                url_or_path=final_path,
                generated_at=export_datetime,
                file_size_bytes=file_size,
            )

        except PermissionError as e:
            error = PresentationError.export_permission_denied()
            error.technical_message = f"File permission error: {str(e)}"
            raise error
        except OSError as e:
            error = PresentationError.export_storage_failed()
            error.technical_message = f"File system error: {str(e)}"
            raise error
        except Exception as e:
            error = PresentationError.export_failed("pdf", str(e))
            raise error
