#!/usr/bin/env python3
"""
Test script for PPTX and PDF export functionality
"""

import sys
import os
from pathlib import Path

# Add the project root to the Python path
project_root = Path(__file__).parent
sys.path.insert(0, str(project_root))


def test_pptx_export():
    """Test PPTX export functionality"""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        # Create a simple test presentation
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Test Presentation"
        subtitle.text = "Generated for testing export functionality"

        # Add a content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)

        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Test Content"

        tf = body_shape.text_frame
        tf.text = "Key points:"

        p = tf.add_paragraph()
        p.text = "First point"
        p.level = 1

        p = tf.add_paragraph()
        p.text = "Second point"
        p.level = 1

        # Save the presentation
        test_file = "test_presentation.pptx"
        prs.save(test_file)

        # Check if file was created
        if os.path.exists(test_file):
            file_size = os.path.getsize(test_file)
            print(f"✅ PPTX export test successful!")
            print(f"   File: {test_file}")
            print(f"   Size: {file_size} bytes")

            # Clean up
            os.remove(test_file)
            return True
        else:
            print("❌ PPTX file was not created")
            return False

    except Exception as e:
        print(f"❌ PPTX export test failed: {e}")
        return False


def test_pdf_export():
    """Test PDF export functionality"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter

        # Create a simple test PDF
        test_file = "test_presentation.pdf"
        c = canvas.Canvas(test_file, pagesize=letter)
        width, height = letter

        # Add title
        c.setFont("Helvetica-Bold", 16)
        c.drawString(72, height - 72, "Test Presentation")

        # Add content
        c.setFont("Helvetica", 12)
        c.drawString(72, height - 108, "Generated for testing export functionality")

        # Add some bullet points
        c.drawString(72, height - 144, "Key points:")
        c.drawString(90, height - 168, "• First point")
        c.drawString(90, height - 192, "• Second point")

        # Save the PDF
        c.save()

        # Check if file was created
        if os.path.exists(test_file):
            file_size = os.path.getsize(test_file)
            print(f"✅ PDF export test successful!")
            print(f"   File: {test_file}")
            print(f"   Size: {file_size} bytes")

            # Clean up
            os.remove(test_file)
            return True
        else:
            print("❌ PDF file was not created")
            return False

    except Exception as e:
        print(f"❌ PDF export test failed: {e}")
        return False


def main():
    """Run all export tests"""
    print("Testing export functionality...")
    print("=" * 50)

    pptx_success = test_pptx_export()
    print()
    pdf_success = test_pdf_export()

    print("=" * 50)
    if pptx_success and pdf_success:
        print("🎉 All export tests passed!")
        return 0
    else:
        print("❌ Some export tests failed")
        return 1


if __name__ == "__main__":
    sys.exit(main())
